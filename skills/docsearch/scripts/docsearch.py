#!/usr/bin/env python3
"""
docsearch - Local semantic document search over ChromaDB.

Indexes any directory into ChromaDB with semantic embeddings via any
OpenAI-compatible embeddings endpoint (LM Studio, Ollama, OpenAI, etc.).
Supports incremental re-sync without duplications using content hashing.

Converts documents using python-docx, python-pptx, pypdfium2, openpyxl
directly (no external CLI dependencies).

Configuration (env vars):
    DOCSEARCH_DIR          Root directory to index (required for index command)
    DOCSEARCH_DB           ChromaDB path (default: ~/.docsearch/chroma_db)
    DOCSEARCH_COLLECTION   Collection name (default: derived from DOCSEARCH_DIR)
    DOCSEARCH_EMBED_URL    Embeddings endpoint (default: http://localhost:1234/v1/embeddings)
    DOCSEARCH_EMBED_MODEL  Embedding model name (default: text-embedding-nomic-embed-text-v1.5)
    DOCSEARCH_EMBED_KEY    API key for embeddings endpoint (default: lmstudio)
    DOCSEARCH_EMBED_DIM    Embedding dimensions (default: 768)

Usage:
    docsearch index [--dir PATH] [--force]
    docsearch search QUERY [--top N] [--filter KEY=VALUE ...]
    docsearch info
    docsearch config
"""

import argparse
import hashlib
import json
import os
import re
import sys
from pathlib import Path
from typing import Optional

import chromadb
import requests

# ---------------------------------------------------------------------------
# Config from environment
# ---------------------------------------------------------------------------

CHROMA_DIR = os.environ.get(
    "DOCSEARCH_DB",
    os.path.expanduser("~/.docsearch/chroma_db"),
)

EMBED_URL = os.environ.get(
    "DOCSEARCH_EMBED_URL",
    "http://localhost:1234/v1/embeddings",
)
EMBED_MODEL = os.environ.get(
    "DOCSEARCH_EMBED_MODEL",
    "text-embedding-nomic-embed-text-v1.5",
)
EMBED_API_KEY = os.environ.get("DOCSEARCH_EMBED_KEY", "lmstudio")
EMBED_DIM = int(os.environ.get("DOCSEARCH_EMBED_DIM", "768"))

# File types we index
INDEXABLE_EXTENSIONS = {
    ".docx", ".pptx", ".pdf", ".xlsx", ".xlsm",
    ".md", ".txt", ".csv",
}

# Max characters per chunk (~500 tokens). Overlap for context continuity.
CHUNK_SIZE = 2000
CHUNK_OVERLAP = 200


def get_default_dir() -> Optional[str]:
    raw = os.environ.get("DOCSEARCH_DIR")
    return os.path.expanduser(raw) if raw else None


def derive_collection_name(root_dir: str) -> str:
    """Derive a stable collection name from the directory path."""
    explicit = os.environ.get("DOCSEARCH_COLLECTION")
    if explicit:
        return explicit
    # Use last path component, sanitized
    name = Path(root_dir).resolve().name
    name = re.sub(r"[^a-zA-Z0-9_-]", "_", name)
    return name or "docsearch_docs"


# ---------------------------------------------------------------------------
# Embedding via OpenAI-compatible endpoint
# ---------------------------------------------------------------------------

def get_embeddings(texts: list[str]) -> list[list[float]]:
    """Get embeddings from an OpenAI-compatible embeddings endpoint."""
    all_embeddings = []
    batch_size = 16
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i + batch_size]
        resp = requests.post(
            EMBED_URL,
            headers={
                "Authorization": f"Bearer {EMBED_API_KEY}",
                "Content-Type": "application/json",
            },
            json={"input": batch, "model": EMBED_MODEL},
            timeout=120,
        )
        resp.raise_for_status()
        data = resp.json()
        embeddings = [item["embedding"] for item in data["data"]]
        all_embeddings.extend(embeddings)
    return all_embeddings


class OpenAICompatibleEmbedding(chromadb.EmbeddingFunction):
    """ChromaDB embedding function backed by any OpenAI-compatible endpoint."""

    def __call__(self, input: list[str]) -> list[list[float]]:
        return get_embeddings(input)


# ---------------------------------------------------------------------------
# Document conversion
# ---------------------------------------------------------------------------

def _convert_docx(file_path: str) -> Optional[str]:
    from docx import Document
    doc = Document(file_path)
    parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n\n".join(parts)


def _convert_pptx(file_path: str) -> Optional[str]:
    from pptx import Presentation
    prs = Presentation(file_path)
    parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_texts.append(" | ".join(cells))
        if slide_texts:
            parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
    return "\n\n".join(parts)


def _convert_pdf(file_path: str) -> Optional[str]:
    import pypdfium2 as pdfium
    pdf = pdfium.PdfDocument(file_path)
    parts = []
    for page in pdf:
        text = page.get_textpage().get_text_range()
        if text and text.strip():
            parts.append(text.strip())
    pdf.close()
    return "\n\n".join(parts)


def _convert_xlsx(file_path: str) -> Optional[str]:
    import openpyxl
    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_parts = [f"[Sheet: {sheet_name}]"]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                sheet_parts.append(" | ".join(cells))
        if len(sheet_parts) > 1:
            parts.append("\n".join(sheet_parts))
    wb.close()
    return "\n\n".join(parts)


def convert_document(file_path: str) -> Optional[str]:
    """Convert a document to plain text. Returns None on failure."""
    ext = Path(file_path).suffix.lower()
    try:
        if ext in (".md", ".txt", ".csv"):
            with open(file_path, "r", errors="replace") as f:
                return f.read()
        elif ext == ".docx":
            return _convert_docx(file_path)
        elif ext == ".pptx":
            return _convert_pptx(file_path)
        elif ext == ".pdf":
            return _convert_pdf(file_path)
        elif ext in (".xlsx", ".xlsm"):
            return _convert_xlsx(file_path)
        else:
            print(f"  WARN: Unsupported format {ext}", file=sys.stderr)
            return None
    except Exception as e:
        print(f"  WARN: Failed to convert {file_path}: {e}", file=sys.stderr)
        return None


# ---------------------------------------------------------------------------
# Chunking
# ---------------------------------------------------------------------------

def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP) -> list[str]:
    """Split text into overlapping chunks."""
    if not text or not text.strip():
        return []

    text = text.strip()
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size

        if end < len(text):
            break_at = text.rfind("\n\n", start + chunk_size // 2, end)
            if break_at == -1:
                break_at = text.rfind(". ", start + chunk_size // 2, end)
                if break_at != -1:
                    break_at += 2
            if break_at > start:
                end = break_at

        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)

        start = end - overlap
        if start >= len(text):
            break

    return chunks


# ---------------------------------------------------------------------------
# Metadata extraction
# ---------------------------------------------------------------------------

def extract_metadata(file_path: str, root_dir: str) -> dict:
    """Extract path-based metadata relative to root_dir."""
    rel = os.path.relpath(file_path, root_dir)
    parts = Path(rel).parts

    meta = {
        "file_path": file_path,
        "relative_path": rel,
        "filename": os.path.basename(file_path),
        "extension": Path(file_path).suffix.lower(),
        "depth_1": parts[0] if len(parts) >= 2 else "",
        "depth_2": parts[1] if len(parts) >= 3 else "",
        "depth_3": parts[2] if len(parts) >= 4 else "",
    }
    return meta


def file_content_hash(file_path: str) -> str:
    h = hashlib.sha256()
    with open(file_path, "rb") as f:
        for block in iter(lambda: f.read(8192), b""):
            h.update(block)
    return h.hexdigest()


# ---------------------------------------------------------------------------
# iCloud helpers
# ---------------------------------------------------------------------------

def is_icloud_placeholder(file_path: str) -> bool:
    return os.path.basename(file_path).startswith(".") and file_path.endswith(".icloud")


def find_icloud_placeholders(root_dir: str) -> list[str]:
    placeholders = []
    for dirpath, _, filenames in os.walk(root_dir):
        for fname in filenames:
            fpath = os.path.join(dirpath, fname)
            if is_icloud_placeholder(fpath):
                placeholders.append(fpath)
    return placeholders


# ---------------------------------------------------------------------------
# Indexing
# ---------------------------------------------------------------------------

def discover_files(root_dir: str) -> list[str]:
    files = []
    for dirpath, dirnames, filenames in os.walk(root_dir):
        dirnames[:] = [d for d in dirnames if not d.startswith(".")]
        for fname in filenames:
            if fname.startswith(".") or fname.startswith("~$") or fname == ".DS_Store":
                continue
            fpath = os.path.join(dirpath, fname)
            ext = Path(fname).suffix.lower()
            if ext in INDEXABLE_EXTENSIONS:
                files.append(fpath)
    return sorted(files)


def do_index(root_dir: str, force: bool = False):
    """Index documents into ChromaDB."""
    collection_name = derive_collection_name(root_dir)
    os.makedirs(CHROMA_DIR, exist_ok=True)

    client = chromadb.PersistentClient(path=CHROMA_DIR)
    embed_fn = OpenAICompatibleEmbedding()

    if force:
        try:
            client.delete_collection(collection_name)
        except Exception:
            pass
        print("Force mode: cleared existing index.")

    collection = client.get_or_create_collection(
        name=collection_name,
        embedding_function=embed_fn,
        metadata={"hnsw:space": "cosine"},
    )

    files = discover_files(root_dir)
    print(f"Found {len(files)} indexable files in {root_dir}")
    print(f"Collection: {collection_name}")

    placeholders = find_icloud_placeholders(root_dir)
    if placeholders:
        print(f"  NOTE: {len(placeholders)} files are in iCloud (not downloaded locally). Skipping these.")
        for p in placeholders[:5]:
            real_name = os.path.basename(p).lstrip(".").replace(".icloud", "")
            print(f"    - {real_name}")
        if len(placeholders) > 5:
            print(f"    ... and {len(placeholders) - 5} more")

    existing_hashes = {}
    if not force and collection.count() > 0:
        all_meta = collection.get(include=["metadatas"])
        for meta in all_meta["metadatas"]:
            fp = meta.get("file_path", "")
            ch = meta.get("content_hash", "")
            if fp and ch:
                existing_hashes[fp] = ch

    if not force and existing_hashes:
        file_set = set(files)
        stale = [fp for fp in existing_hashes if fp not in file_set]
        if stale:
            print(f"  Removing {len(stale)} stale entries (files deleted from disk)...")
            for fp in stale:
                old_ids = collection.get(where={"file_path": fp}, include=[])
                if old_ids["ids"]:
                    collection.delete(ids=old_ids["ids"])

    indexed = 0
    skipped = 0
    errors = 0

    for i, fpath in enumerate(files, 1):
        fname = os.path.basename(fpath)
        print(f"[{i}/{len(files)}] {fname}", end=" ... ")

        try:
            chash = file_content_hash(fpath)

            if not force and fpath in existing_hashes and existing_hashes[fpath] == chash:
                print("unchanged, skipping")
                skipped += 1
                continue

            if fpath in existing_hashes:
                old_ids = collection.get(where={"file_path": fpath}, include=[])
                if old_ids["ids"]:
                    collection.delete(ids=old_ids["ids"])
                    print(f"(removed {len(old_ids['ids'])} old chunks) ", end="")

            text = convert_document(fpath)
            if not text or not text.strip():
                print("empty/failed")
                errors += 1
                continue

            chunks = chunk_text(text)
            if not chunks:
                print("no chunks")
                errors += 1
                continue

            meta = extract_metadata(fpath, root_dir)
            mtime = os.path.getmtime(fpath)
            fsize = os.path.getsize(fpath)

            ids = []
            documents = []
            metadatas = []

            for ci, chunk in enumerate(chunks):
                chunk_id = f"{chash}_{ci}"
                ids.append(chunk_id)
                documents.append(chunk)
                metadatas.append({
                    "file_path": meta["file_path"],
                    "relative_path": meta["relative_path"],
                    "filename": meta["filename"],
                    "extension": meta["extension"],
                    "depth_1": meta["depth_1"],
                    "depth_2": meta["depth_2"],
                    "depth_3": meta["depth_3"],
                    "content_hash": chash,
                    "chunk_index": ci,
                    "total_chunks": len(chunks),
                    "modified_time": mtime,
                    "file_size": fsize,
                })

            collection.upsert(
                ids=ids,
                documents=documents,
                metadatas=metadatas,
            )

            indexed += 1
            print(f"{len(chunks)} chunks")

        except Exception as e:
            print(f"ERROR: {e}")
            errors += 1

    print(f"\nDone. Indexed: {indexed}, Skipped (unchanged): {skipped}, Errors: {errors}")
    print(f"Total chunks in collection: {collection.count()}")


# ---------------------------------------------------------------------------
# Search
# ---------------------------------------------------------------------------

def do_search(
    query: str,
    top_n: int = 5,
    filters: Optional[dict] = None,
):
    """Semantic search over indexed documents."""
    default_dir = get_default_dir()
    if not default_dir:
        print("Set DOCSEARCH_DIR to the directory you indexed.", file=sys.stderr)
        sys.exit(1)

    collection_name = derive_collection_name(default_dir)
    client = chromadb.PersistentClient(path=CHROMA_DIR)
    embed_fn = OpenAICompatibleEmbedding()

    try:
        collection = client.get_collection(
            name=collection_name,
            embedding_function=embed_fn,
        )
    except Exception:
        print(f"No index found for collection '{collection_name}'. Run 'docsearch index' first.", file=sys.stderr)
        sys.exit(1)

    where = None
    if filters:
        conditions = [{"$eq": {k: v}} for k, v in filters.items()]
        if len(conditions) == 1:
            where = conditions[0]
        elif len(conditions) > 1:
            where = {"$and": conditions}

    results = collection.query(
        query_texts=[query],
        n_results=min(top_n * 3, max(collection.count(), 1)),
        where=where,
        include=["documents", "metadatas", "distances"],
    )

    if not results["ids"][0]:
        print("No results found.")
        return

    seen_files = {}
    for idx in range(len(results["ids"][0])):
        meta = results["metadatas"][0][idx]
        distance = results["distances"][0][idx]
        document = results["documents"][0][idx]
        fpath = meta["file_path"]

        if fpath not in seen_files:
            seen_files[fpath] = {
                "distance": distance,
                "meta": meta,
                "preview": document[:300],
            }

    ranked = sorted(seen_files.items(), key=lambda x: x[1]["distance"])[:top_n]

    output = []
    for rank, (fpath, info) in enumerate(ranked, 1):
        meta = info["meta"]
        score = 1.0 - info["distance"]
        entry = {
            "rank": rank,
            "score": round(score, 4),
            "file": fpath,
            "filename": meta["filename"],
            "relative_path": meta["relative_path"],
            "type": meta["extension"],
            "depth_1": meta.get("depth_1", ""),
            "depth_2": meta.get("depth_2", ""),
            "preview": info["preview"].replace("\n", " ")[:200],
        }
        output.append(entry)

    if sys.stdout.isatty():
        for entry in output:
            path_hint = "/".join(filter(None, [entry["depth_1"], entry["depth_2"]]))
            path_str = f" [{path_hint}]" if path_hint else ""
            print(f"\n  #{entry['rank']}  {entry['filename']}{path_str}")
            print(f"      Score: {entry['score']}  |  Type: {entry['type']}")
            print(f"      Path:  {entry['file']}")
            print(f"      Preview: {entry['preview'][:120]}...")
        print()
    else:
        print(json.dumps(output, indent=2))


# ---------------------------------------------------------------------------
# Info
# ---------------------------------------------------------------------------

def do_info():
    default_dir = get_default_dir()

    if not os.path.exists(CHROMA_DIR):
        print("No index found. Run 'docsearch index' first.")
        return

    client = chromadb.PersistentClient(path=CHROMA_DIR)

    # List all collections if no DOCSEARCH_DIR set
    if not default_dir:
        collections = client.list_collections()
        print(f"ChromaDB path: {CHROMA_DIR}")
        print(f"Collections ({len(collections)}):")
        for c in collections:
            col = client.get_collection(c.name)
            print(f"  {c.name}: {col.count()} chunks")
        return

    collection_name = derive_collection_name(default_dir)
    try:
        collection = client.get_collection(name=collection_name)
    except Exception:
        print(f"No collection '{collection_name}'. Run 'docsearch index' first.")
        return

    count = collection.count()
    all_meta = collection.get(include=["metadatas"])

    files = set()
    extensions = {}
    depth1_counts = {}

    for meta in all_meta["metadatas"]:
        files.add(meta.get("file_path", ""))
        ext = meta.get("extension", "")
        if ext:
            extensions[ext] = extensions.get(ext, 0) + 1
        d1 = meta.get("depth_1", "")
        if d1:
            depth1_counts[d1] = depth1_counts.get(d1, 0) + 1

    print(f"ChromaDB path: {CHROMA_DIR}")
    print(f"Collection:    {collection_name}")
    print(f"Root dir:      {default_dir}")
    print(f"Total chunks:  {count}")
    print(f"Total files:   {len(files)}")
    if depth1_counts:
        print(f"\nBy top-level directory:")
        for d, cnt in sorted(depth1_counts.items(), key=lambda x: -x[1]):
            print(f"  {d}: {cnt} chunks")
    print(f"\nBy file type:")
    for ext, cnt in sorted(extensions.items(), key=lambda x: -x[1]):
        print(f"  {ext}: {cnt} chunks")


def do_config():
    """Print current configuration."""
    default_dir = get_default_dir()
    collection_name = derive_collection_name(default_dir) if default_dir else "(set DOCSEARCH_DIR first)"
    print(f"DOCSEARCH_DIR        = {default_dir or '(not set)'}")
    print(f"DOCSEARCH_DB         = {CHROMA_DIR}")
    print(f"DOCSEARCH_COLLECTION = {collection_name}")
    print(f"DOCSEARCH_EMBED_URL  = {EMBED_URL}")
    print(f"DOCSEARCH_EMBED_MODEL= {EMBED_MODEL}")
    print(f"DOCSEARCH_EMBED_DIM  = {EMBED_DIM}")


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Local semantic document search with ChromaDB",
        prog="docsearch",
    )
    sub = parser.add_subparsers(dest="command")

    # index
    idx = sub.add_parser("index", help="Index documents into ChromaDB")
    default_dir = get_default_dir()
    idx.add_argument(
        "--dir",
        default=default_dir,
        required=(default_dir is None),
        help="Root directory to index (or set DOCSEARCH_DIR env var)",
    )
    idx.add_argument("--force", action="store_true", help="Re-index everything (ignore hashes)")

    # search
    srch = sub.add_parser("search", help="Semantic search over indexed documents")
    srch.add_argument("query", nargs="+", help="Search query")
    srch.add_argument("--top", type=int, default=5, help="Number of results")
    srch.add_argument(
        "--filter",
        action="append",
        metavar="KEY=VALUE",
        help="Filter by metadata field (e.g. --filter extension=.pdf --filter depth_1=customers)",
    )

    # info
    sub.add_parser("info", help="Show index statistics")

    # config
    sub.add_parser("config", help="Show current configuration")

    args = parser.parse_args()

    if args.command == "index":
        do_index(args.dir, force=args.force)
    elif args.command == "search":
        query = " ".join(args.query)
        filters = {}
        if args.filter:
            for f in args.filter:
                if "=" in f:
                    k, v = f.split("=", 1)
                    filters[k.strip()] = v.strip()
        do_search(query, top_n=args.top, filters=filters or None)
    elif args.command == "info":
        do_info()
    elif args.command == "config":
        do_config()
    else:
        parser.print_help()


if __name__ == "__main__":
    main()
